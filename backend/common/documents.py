# SPDX-License-Identifier: Apache-2.0
# Copyright 2026 Universität Osnabrück (virtUOS)

"""Extract plain text from uploaded documents for the AI features.

Starts with PDF (pypdf); presentation formats (PPTX/ODP) are added with the
"generate questions from documents" feature. Text is collapsed and truncated
so the prompt (and cost) stays bounded. No OCR — image-only files raise
``DocumentTextError``."""
import re

MAX_CHARS = 12000


class DocumentTextError(Exception):
    """The document could not be turned into usable text."""


def _collapse(parts, max_chars):
    text = re.sub(r"\s+", " ", " ".join(parts)).strip()
    if not text:
        raise DocumentTextError("Kein extrahierbarer Text (evtl. gescanntes PDF).")
    return text[:max_chars]


def extract_pdf_text(file, *, max_chars=MAX_CHARS):
    from pypdf import PdfReader

    try:
        reader = PdfReader(file)
        parts = [page.extract_text() or "" for page in reader.pages]
    except Exception as exc:  # pypdf raises a variety of errors on bad input
        raise DocumentTextError(str(exc)) from exc
    return _collapse(parts, max_chars)


def extract_pptx_text(file, *, max_chars=MAX_CHARS):
    from pptx import Presentation

    try:
        presentation = Presentation(file)
        parts = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        parts.extend(cell.text for cell in row.cells)
    except Exception as exc:
        raise DocumentTextError(str(exc)) from exc
    return _collapse(parts, max_chars)


def extract_odp_text(file, *, max_chars=MAX_CHARS):
    from odf import teletype
    from odf.opendocument import load
    from odf.text import P

    try:
        document = load(file)
        parts = [teletype.extractText(p) for p in document.getElementsByType(P)]
    except Exception as exc:
        raise DocumentTextError(str(exc)) from exc
    return _collapse(parts, max_chars)


_EXTRACTORS = {
    ".pdf": extract_pdf_text,
    ".pptx": extract_pptx_text,
    ".odp": extract_odp_text,
}


def extract_text(file, filename, *, max_chars=MAX_CHARS):
    """Dispatch on the file extension. Raises DocumentTextError for an
    unsupported type or an unreadable/scanned file."""
    name = (filename or "").lower()
    for suffix, extractor in _EXTRACTORS.items():
        if name.endswith(suffix):
            return extractor(file, max_chars=max_chars)
    raise DocumentTextError(
        "Nicht unterstütztes Format. Erlaubt sind PDF, PPTX und ODP."
    )
