# SPDX-License-Identifier: Apache-2.0
# Copyright 2026 Universität Osnabrück (virtUOS)

from unittest.mock import MagicMock, patch
from urllib import error

from django.contrib.auth import get_user_model
from django.test import TestCase, override_settings
from django.utils import translation

from .markdown import render_markdown
from .models import Page, SiteConfig

User = get_user_model()

def _mock_response(body):
    """A urlopen() context manager whose .read() yields the given JSON body."""
    cm = MagicMock()
    cm.__enter__.return_value.read.return_value = body.encode("utf-8")
    return cm


class PublicSiteTests(TestCase):
    """Branding, footer pages, page detail and the data registry are public."""

    def test_seeded_legal_pages_are_drafts(self):
        # The seed migration creates them unpublished, so not in the footer.
        self.assertTrue(Page.objects.filter(slug="impressum").exists())
        self.assertTrue(Page.objects.filter(slug="datenschutz").exists())
        footer = self.client.get("/api/pages/").json()
        self.assertEqual(footer, [])

    def test_published_page_shows_in_footer_and_detail(self):
        # Outside a request, the active language is LANGUAGE_CODE ("en"), not
        # the content-canonical MODELTRANSLATION_DEFAULT_LANGUAGE ("de") —
        # override so this fixture's title/body land in the _de columns like
        # real authored content would (#33 MR2).
        with translation.override("de"):
            Page.objects.create(slug="ueber-uns", title="Über uns", body="# Hallo")
        footer = self.client.get("/api/pages/").json()
        self.assertEqual([p["slug"] for p in footer], ["ueber-uns"])
        # title/body are {"de", "en"} maps (#33 MR2), resolved client-side.
        self.assertEqual(footer[0]["title"], {"de": "Über uns", "en": ""})
        detail = self.client.get("/api/pages/ueber-uns/").json()
        self.assertEqual(detail["title"], {"de": "Über uns", "en": ""})
        self.assertEqual(detail["body"], {"de": "# Hallo", "en": ""})

    def test_unpublished_page_detail_404(self):
        Page.objects.create(slug="entwurf", title="Entwurf", is_published=False)
        self.assertEqual(self.client.get("/api/pages/entwurf/").status_code, 404)

    def test_site_returns_landing_and_logo(self):
        SiteConfig.load()  # ensure singleton
        payload = self.client.get("/api/site/").json()
        self.assertIn("landing_text", payload)
        # landing_text/closing_info are {"de", "en"} maps (#33 MR2).
        self.assertEqual(payload["landing_text"], {"de": "", "en": ""})
        self.assertIsNone(payload["logo"])

    def test_data_collection_registry(self):
        payload = self.client.get("/api/data-collection/").json()
        self.assertTrue(len(payload["collected"]) >= 1)
        self.assertIn("category", payload["collected"][0])
        self.assertTrue(len(payload["not_collected"]) >= 1)


class MarkdownRenderTests(TestCase):
    """Server-side Markdown for the participant closing screen (#24)."""

    def test_links_and_images_survive_scripts_stripped(self):
        html = render_markdown(
            "**Hallo** [Link](https://e.com) ![Bild](https://e.com/i.png)\n\n"
            "<script>alert(1)</script>"
        )
        self.assertIn("<strong>Hallo</strong>", html)
        self.assertIn('<img', html)
        self.assertIn('href="https://e.com"', html)
        # Shared allowlist (#49) forces rel="noopener" only (not "noopener
        # noreferrer" as the old standalone markdown allowlist did).
        self.assertIn('rel="noopener"', html)
        self.assertNotIn("<script", html)

    def test_empty_is_empty(self):
        self.assertEqual(render_markdown(""), "")


class ManageSiteTests(TestCase):
    def setUp(self):
        self.admin = User.objects.create_user(username="chef", is_staff=True)
        self.plain = User.objects.create_user(username="hans")

    def test_manage_requires_staff(self):
        self.assertEqual(self.client.get("/api/manage/site/").status_code, 403)
        self.client.force_login(self.plain)
        self.assertEqual(self.client.get("/api/manage/site/").status_code, 403)
        self.client.force_login(self.admin)
        self.assertEqual(self.client.get("/api/manage/site/").status_code, 200)

    def test_admin_edits_landing_text(self):
        self.client.force_login(self.admin)
        # A plain (legacy) string is written to the canonical language only.
        response = self.client.put(
            "/api/manage/site/", {"landing_text": "Willkommen!"},
            content_type="application/json",
        )
        self.assertEqual(response.status_code, 200)
        self.assertEqual(SiteConfig.load().landing_text_de, "Willkommen!")
        self.assertEqual(response.json()["landing_text"], {"de": "Willkommen!", "en": ""})

    def test_admin_edits_ai_notice_bilingually_and_public_reads_it(self):
        self.client.force_login(self.admin)
        response = self.client.put(
            "/api/manage/site/",
            {
                "ai_notice": {"de": "Externes Modell.", "en": "External model."},
                "ai_notice_url": "https://uni.example/datenschutz",
            },
            content_type="application/json",
        )
        self.assertEqual(response.status_code, 200)
        config = SiteConfig.load()
        self.assertEqual(config.ai_notice_de, "Externes Modell.")
        self.assertEqual(config.ai_notice_en, "External model.")
        self.assertEqual(config.ai_notice_url, "https://uni.example/datenschutz")
        # The public site endpoint exposes it (the SPA reads it there).
        self.client.logout()
        public = self.client.get("/api/site/").json()
        self.assertEqual(public["ai_notice"], {"de": "Externes Modell.", "en": "External model."})
        self.assertEqual(public["ai_notice_url"], "https://uni.example/datenschutz")

    def test_ai_notice_empty_by_default(self):
        self.assertEqual(self.client.get("/api/site/").json()["ai_notice"], {"de": "", "en": ""})

    def test_ai_notice_can_link_an_internal_page_by_slug(self):
        Page.objects.get_or_create(slug="ds-eigen", defaults={"title": "Datenschutz"})
        self.client.force_login(self.admin)
        response = self.client.put(
            "/api/manage/site/",
            {"ai_notice_page": "ds-eigen"},
            content_type="application/json",
        )
        self.assertEqual(response.status_code, 200)
        self.assertEqual(SiteConfig.load().ai_notice_page.slug, "ds-eigen")
        # Exposed as the slug on the public endpoint (the banner links to it).
        self.assertEqual(self.client.get("/api/site/").json()["ai_notice_page"], "ds-eigen")

    def test_ai_notice_page_null_by_default(self):
        self.assertIsNone(self.client.get("/api/site/").json()["ai_notice_page"])

    def test_admin_edits_landing_text_via_map(self):
        self.client.force_login(self.admin)
        response = self.client.put(
            "/api/manage/site/",
            {"landing_text": {"de": "Willkommen!", "en": "Welcome!"}},
            content_type="application/json",
        )
        self.assertEqual(response.status_code, 200)
        config = SiteConfig.load()
        self.assertEqual(config.landing_text_de, "Willkommen!")
        self.assertEqual(config.landing_text_en, "Welcome!")
        self.assertEqual(
            response.json()["landing_text"], {"de": "Willkommen!", "en": "Welcome!"}
        )

    def test_admin_edits_closing_info(self):
        self.client.force_login(self.admin)
        response = self.client.put(
            "/api/manage/site/",
            {"landing_text": "Hi", "closing_info": "# Danke\n[Feedback](https://e.com)"},
            content_type="application/json",
        )
        self.assertEqual(response.status_code, 200)
        self.assertIn("Danke", SiteConfig.load().closing_info_de)

    def test_admin_edits_landing_text_only_leaves_closing_info_untouched(self):
        self.client.force_login(self.admin)
        self.client.put(
            "/api/manage/site/",
            {"landing_text": "Hi", "closing_info": "Bleib"},
            content_type="application/json",
        )
        response = self.client.put(
            "/api/manage/site/", {"landing_text": "Neu"},
            content_type="application/json",
        )
        self.assertEqual(response.status_code, 200)
        self.assertEqual(SiteConfig.load().closing_info_de, "Bleib")

    def test_admin_page_crud_and_reorder(self):
        self.client.force_login(self.admin)
        a = self.client.post(
            "/api/manage/pages/", {"slug": "a", "title": "A"},
            content_type="application/json",
        ).json()
        b = self.client.post(
            "/api/manage/pages/", {"slug": "b", "title": "B"},
            content_type="application/json",
        ).json()
        # New pages append to the end.
        self.assertLess(a["footer_order"], b["footer_order"])
        # Reorder swaps them.
        self.client.post(
            "/api/manage/pages/reorder/", {"order": [b["id"], a["id"]]},
            content_type="application/json",
        )
        self.assertEqual(Page.objects.get(pk=b["id"]).footer_order, 0)
        self.assertEqual(Page.objects.get(pk=a["id"]).footer_order, 1)

    def test_plain_user_cannot_create_page(self):
        self.client.force_login(self.plain)
        response = self.client.post(
            "/api/manage/pages/", {"slug": "x", "title": "X"},
            content_type="application/json",
        )
        self.assertEqual(response.status_code, 403)

    def test_landing_text_html_sanitized_on_write(self):
        # The editor now sends HTML (#49); validate_landing_text must run
        # it through the shared allowlist, stripping <script>.
        self.client.force_login(self.admin)
        response = self.client.put(
            "/api/manage/site/",
            {"landing_text": {"de": '<p>hi<script>x()</script></p>', "en": ""}},
            content_type="application/json",
        )
        self.assertEqual(response.status_code, 200)
        self.assertEqual(SiteConfig.load().landing_text_de, "<p>hi</p>")

    def test_page_body_html_sanitized_on_write(self):
        # h1 is not in the shared allowlist and gets unwrapped to text (see
        # HtmlSanitizeTests.test_allows_headings_h2_h3_but_drops_h1); h2 survives.
        self.client.force_login(self.admin)
        response = self.client.post(
            "/api/manage/pages/",
            {"slug": "x", "title": "T", "body": {"de": "<h1>no</h1><h2>yes</h2>", "en": ""}},
            content_type="application/json",
        )
        self.assertEqual(response.status_code, 201)
        page = Page.objects.get(slug="x")
        self.assertNotIn("<h1", page.body_de)
        self.assertIn("<h2>yes</h2>", page.body_de)


class DocumentExtractionTests(TestCase):
    """Text extraction + format dispatch for the document AI feature."""

    def test_pptx_extraction(self):
        import io

        from pptx import Presentation
        from pptx.util import Inches

        from common.documents import extract_text

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        box.text_frame.text = "Photosynthese Grundlagen"
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        self.assertIn("Photosynthese", extract_text(buf, "folien.pptx"))

    def test_unsupported_extension_raises(self):
        import io

        from common.documents import DocumentTextError, extract_text

        with self.assertRaises(DocumentTextError):
            extract_text(io.BytesIO(b"plain"), "notes.txt")

    def test_empty_document_raises(self):
        import io

        from pptx import Presentation

        from common.documents import DocumentTextError, extract_text

        prs = Presentation()  # no slides, no text
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        with self.assertRaises(DocumentTextError):
            extract_text(buf, "leer.pptx")


LT_ON = {"CONTENT_TRANSLATION_PROVIDER": "libretranslate", "LIBRETRANSLATE_URL": "http://lt"}
LT_OFF = {"CONTENT_TRANSLATION_PROVIDER": "none", "LIBRETRANSLATE_URL": ""}


class TranslateEndpointTests(TestCase):
    """POST /api/translate/ (#33 MR2)."""

    def setUp(self):
        self.user = User.objects.create_user(username="frank")

    def test_unauthenticated_forbidden(self):
        response = self.client.post(
            "/api/translate/",
            {"text": "Hallo", "source": "de", "target": "en"},
            content_type="application/json",
        )
        self.assertEqual(response.status_code, 403)

    @override_settings(**LT_OFF)
    def test_disabled_returns_503(self):
        self.client.force_login(self.user)
        response = self.client.post(
            "/api/translate/",
            {"text": "Hallo", "source": "de", "target": "en"},
            content_type="application/json",
        )
        self.assertEqual(response.status_code, 503)

    @override_settings(**LT_ON)
    def test_success_returns_translated_text(self):
        self.client.force_login(self.user)
        body = '{"translatedText": "Hello"}'
        with patch(
            "basicbar_integrations.translation_service.request.urlopen",
            return_value=_mock_response(body),
        ):
            response = self.client.post(
                "/api/translate/",
                {"text": "Hallo", "source": "de", "target": "en"},
                content_type="application/json",
            )
        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json(), {"translated": "Hello"})

    @override_settings(**LT_ON)
    def test_upstream_error_returns_502_not_stacktrace(self):
        self.client.force_login(self.user)
        with patch(
            "basicbar_integrations.translation_service.request.urlopen",
            side_effect=error.URLError("boom"),
        ):
            response = self.client.post(
                "/api/translate/",
                {"text": "Hallo", "source": "de", "target": "en"},
                content_type="application/json",
            )
        self.assertEqual(response.status_code, 502)
        self.assertIn("detail", response.json())

    @override_settings(**LT_ON)
    def test_html_result_is_sanitized(self):
        self.client.force_login(self.user)
        body = '{"translatedText": "<p>Hi</p><script>alert(1)</script>"}'
        with patch(
            "basicbar_integrations.translation_service.request.urlopen",
            return_value=_mock_response(body),
        ):
            response = self.client.post(
                "/api/translate/",
                {"text": "<p>Hallo</p>", "source": "de", "target": "en", "format": "html"},
                content_type="application/json",
            )
        self.assertEqual(response.status_code, 200)
        translated = response.json()["translated"]
        self.assertIn("<p>Hi</p>", translated)
        self.assertNotIn("<script", translated)

    @override_settings(**LT_ON)
    def test_unknown_language_rejected(self):
        self.client.force_login(self.user)
        response = self.client.post(
            "/api/translate/",
            {"text": "Hallo", "source": "de", "target": "fr"},
            content_type="application/json",
        )
        self.assertEqual(response.status_code, 400)

    @override_settings(**LT_ON)
    def test_same_language_returns_text_unchanged(self):
        self.client.force_login(self.user)
        # No urlopen mock needed — translate() short-circuits before any call.
        response = self.client.post(
            "/api/translate/",
            {"text": "Hallo", "source": "de", "target": "de"},
            content_type="application/json",
        )
        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json(), {"translated": "Hallo"})

    @override_settings(**LT_ON)
    def test_blank_text_returns_unchanged(self):
        self.client.force_login(self.user)
        response = self.client.post(
            "/api/translate/",
            {"text": "", "source": "de", "target": "en"},
            content_type="application/json",
        )
        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json(), {"translated": ""})


class RenderMarkdownAllowlistTests(TestCase):
    def test_markdown_link_gets_rel_noopener(self):
        out = render_markdown("[x](https://example.org)")
        self.assertIn('rel="noopener"', out)
        self.assertIn('href="https://example.org"', out)

    def test_markdown_heading_normalizes_to_allowed_subset(self):
        # '#' -> <h1>, which is not in the shared allowlist -> unwrapped to text.
        out = render_markdown("# Title\n\nBody")
        self.assertNotIn("<h1", out)
        self.assertIn("Title", out)

    def test_markdown_list_survives(self):
        out = render_markdown("- a\n- b")
        self.assertIn("<ul>", out)
        self.assertIn("<li>a</li>", out)

    def test_markdown_empty_returns_empty(self):
        self.assertEqual(render_markdown(""), "")
